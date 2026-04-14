from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT = Path(__file__).resolve().parent / "etat_avancement_projet_2026-04-01.pptx"


COLORS = {
    "navy": RGBColor(11, 31, 58),
    "slate": RGBColor(51, 78, 104),
    "teal": RGBColor(27, 153, 139),
    "gold": RGBColor(244, 185, 66),
    "coral": RGBColor(229, 107, 111),
    "ink": RGBColor(28, 37, 44),
    "muted": RGBColor(102, 117, 127),
    "bg": RGBColor(247, 245, 242),
    "white": RGBColor(255, 255, 255),
    "line": RGBColor(219, 226, 232),
    "green": RGBColor(47, 158, 68),
    "amber": RGBColor(240, 140, 0),
    "red": RGBColor(224, 49, 49),
}


def set_background(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_top_band(slide):
    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.45)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = COLORS["navy"]
    band.line.fill.background()


def add_footer(slide, text="Source: audit du depot local au 1 avril 2026"):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(6.95), Inches(12), Inches(0.25))
    p = box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(10)
    run.font.color.rgb = COLORS["muted"]


def add_title(slide, title, subtitle=None):
    add_top_band(slide)
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.7), Inches(9.8), Inches(0.75))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Aptos Display"
    run.font.bold = True
    run.font.size = Pt(28)
    run.font.color.rgb = COLORS["navy"]
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.72), Inches(1.42), Inches(9.5), Inches(0.45))
        sp = sub_box.text_frame.paragraphs[0]
        srun = sp.add_run()
        srun.text = subtitle
        srun.font.name = "Aptos"
        srun.font.size = Pt(14)
        srun.font.color.rgb = COLORS["slate"]


def add_text(slide, x, y, w, h, text, size=18, color="ink", bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = COLORS[color]
    return box


def add_bullets(slide, x, y, w, h, bullets, size=18, color="ink"):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for idx, item in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = COLORS[color]
        p.space_after = Pt(10)
        p.bullet = True
    return box


def add_card(slide, x, y, w, h, title, body, accent):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = COLORS["white"]
    card.line.color.rgb = COLORS["line"]

    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, Inches(0.16), h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS[accent]
    bar.line.fill.background()

    add_text(slide, x + Inches(0.28), y + Inches(0.18), w - Inches(0.4), Inches(0.35), title, 18, "navy", True)
    add_text(slide, x + Inches(0.28), y + Inches(0.62), w - Inches(0.4), h - Inches(0.8), body, 13, "slate")


def add_status_pill(slide, x, y, w, label, accent):
    pill = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, Inches(0.34))
    pill.fill.solid()
    pill.fill.fore_color.rgb = COLORS[accent]
    pill.line.fill.background()
    tf = pill.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.name = "Aptos"
    run.font.bold = True
    run.font.size = Pt(11)
    run.font.color.rgb = COLORS["white"]


def add_timeline_node(slide, x, title, date, bullets, accent):
    circle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x, Inches(2.25), Inches(0.42), Inches(0.42))
    circle.fill.solid()
    circle.fill.fore_color.rgb = COLORS[accent]
    circle.line.fill.background()

    add_text(slide, x - Inches(0.2), Inches(1.65), Inches(1.0), Inches(0.35), date, 12, "muted", True, PP_ALIGN.CENTER)
    add_text(slide, x - Inches(0.45), Inches(2.8), Inches(1.5), Inches(0.45), title, 15, "navy", True, PP_ALIGN.CENTER)
    add_bullets(slide, x - Inches(0.75), Inches(3.18), Inches(2.3), Inches(2.1), bullets, 12, "slate")


def add_arch_box(slide, x, y, w, h, title, subtitle, accent):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["white"]
    shape.line.color.rgb = COLORS[accent]
    add_text(slide, x + Inches(0.16), y + Inches(0.15), w - Inches(0.32), Inches(0.3), title, 18, "navy", True)
    add_text(slide, x + Inches(0.16), y + Inches(0.55), w - Inches(0.32), h - Inches(0.65), subtitle, 12, "slate")
    return shape


def add_arrow(slide, x, y, w):
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, x, y, w, Inches(0.28))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["gold"]
    line.line.fill.background()


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    layout = prs.slide_layouts[6]

    # Slide 1
    slide = prs.slides.add_slide(layout)
    set_background(slide, COLORS["bg"])
    hero = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(0.85), Inches(12.2), Inches(5.4)
    )
    hero.fill.solid()
    hero.fill.fore_color.rgb = COLORS["white"]
    hero.line.color.rgb = COLORS["line"]

    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(0.85), Inches(0.24), Inches(5.4)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = COLORS["teal"]
    accent.line.fill.background()

    add_text(
        slide,
        Inches(1.0),
        Inches(1.2),
        Inches(8.6),
        Inches(1.2),
        "Etat d'avancement du projet ATM ISO 8583 Gateway",
        30,
        "navy",
        True,
    )
    add_text(
        slide,
        Inches(1.0),
        Inches(2.25),
        Inches(8.8),
        Inches(0.8),
        "Presentation de suivi basee sur le depot local, les tests et les artefacts de build verifies le 1 avril 2026.",
        18,
        "slate",
    )
    add_status_pill(slide, Inches(1.0), Inches(3.25), Inches(2.0), "Backend valide", "green")
    add_status_pill(slide, Inches(3.15), Inches(3.25), Inches(2.0), "Frontend build OK", "teal")
    add_status_pill(slide, Inches(5.3), Inches(3.25), Inches(2.5), "Docker compose pret", "gold")
    add_status_pill(slide, Inches(7.95), Inches(3.25), Inches(2.8), "Industrialisation a finir", "coral")

    add_bullets(
        slide,
        Inches(1.0),
        Inches(4.1),
        Inches(10.8),
        Inches(1.7),
        [
            "Positionnement: passerelle Spring Boot qui traduit JSON/XML vers ISO 8583 et dialogue avec un switch via TCP.",
            "Constat global: la base fonctionnelle est solide et demonstrable, mais la branche reste en phase de consolidation.",
        ],
        16,
        "ink",
    )
    add_footer(slide)

    # Slide 2
    slide = prs.slides.add_slide(layout)
    set_background(slide, COLORS["bg"])
    add_title(slide, "Synthese Executive", "Lecture rapide de l'etat du projet")
    add_card(
        slide,
        Inches(0.8),
        Inches(1.8),
        Inches(2.8),
        Inches(1.65),
        "Noyau ISO 8583",
        "Controleur REST, codec jPOS, canal TCP, echo reseau et mapping reponse en place.",
        "green",
    )
    add_card(
        slide,
        Inches(3.75),
        Inches(1.8),
        Inches(2.8),
        Inches(1.65),
        "Observabilite",
        "Metrics, evenements recents, erreurs recentes et dashboard Angular integre sous /dashboard/.",
        "teal",
    )
    add_card(
        slide,
        Inches(6.7),
        Inches(1.8),
        Inches(2.8),
        Inches(1.65),
        "Interop XML",
        "Facade powerCARD /api/powercard/direct-debit mappee vers ISO 8583 MTI 1200.",
        "gold",
    )
    add_card(
        slide,
        Inches(9.65),
        Inches(1.8),
        Inches(2.8),
        Inches(1.65),
        "Deploiement",
        "Dockerfile et docker-compose valides syntaxiquement; execution complete non rejouee dans cette session.",
        "coral",
    )
    add_bullets(
        slide,
        Inches(0.95),
        Inches(4.05),
        Inches(11.7),
        Inches(2.3),
        [
            "Verdict: projet pret pour une demonstration technique encadree.",
            "Frein principal avant une livraison plus formelle: beaucoup de changements locaux non figes (39 modifies, 35 nouveaux, 3 suppressions).",
            "La documentation publique du depot est en retard par rapport au code actuel.",
        ],
        18,
        "ink",
    )
    add_footer(slide)

    # Slide 3
    slide = prs.slides.add_slide(layout)
    set_background(slide, COLORS["bg"])
    add_title(slide, "Périmètre Livré", "Fonctionnalites presentes dans le code au 1 avril 2026")
    add_card(
        slide,
        Inches(0.75),
        Inches(1.9),
        Inches(5.9),
        Inches(3.9),
        "Backend et integration",
        "Endpoints /api/iso8583/send, /authorize, /financial, /presentment, /reversal, /echo, /health et /status.\n\n"
        "Validation des requetes, conversion JSON -> ISOMsg -> JSON, resolution du statut HTTP, gestion globale des erreurs et tracage X-Request-ID.\n\n"
        "Simulation locale avec mock switch et documentation OpenAPI / Swagger.",
        "green",
    )
    add_card(
        slide,
        Inches(6.85),
        Inches(1.9),
        Inches(5.7),
        Inches(3.9),
        "Supervision et interface",
        "Service de monitoring en memoire avec total transactions, succes, declins, erreurs, latence moyenne, P95 et historique recent.\n\n"
        "Frontend Angular compile et copie dans Spring Boot pour exposition sous /dashboard/.\n\n"
        "Facade XML powerCARD pour virements direct-debit.",
        "teal",
    )
    add_footer(slide)

    # Slide 4
    slide = prs.slides.add_slide(layout)
    set_background(slide, COLORS["bg"])
    add_title(slide, "Architecture Actuelle", "Vue simplifiee de la chaine de traitement")
    add_arch_box(
        slide,
        Inches(0.7),
        Inches(2.25),
        Inches(2.2),
        Inches(1.3),
        "Clients",
        "REST JSON\nXML powerCARD",
        "teal",
    )
    add_arrow(slide, Inches(3.02), Inches(2.74), Inches(0.6))
    add_arch_box(
        slide,
        Inches(3.55),
        Inches(2.05),
        Inches(2.5),
        Inches(1.7),
        "Gateway Spring Boot",
        "Controllers\nValidation\nGlobalExceptionHandler",
        "navy",
    )
    add_arrow(slide, Inches(6.15), Inches(2.74), Inches(0.6))
    add_arch_box(
        slide,
        Inches(6.7),
        Inches(2.05),
        Inches(2.5),
        Inches(1.7),
        "Service ISO 8583",
        "Codec jPOS\nPackager XML\nStatus resolver",
        "gold",
    )
    add_arrow(slide, Inches(9.3), Inches(2.74), Inches(0.6))
    add_arch_box(
        slide,
        Inches(9.85),
        Inches(2.25),
        Inches(2.2),
        Inches(1.3),
        "Reseau TCP",
        "Iso8583Channel\nMock switch / host",
        "coral",
    )
    add_arch_box(
        slide,
        Inches(3.1),
        Inches(4.55),
        Inches(3.2),
        Inches(1.2),
        "MonitoringService",
        "Compteurs, latence, historique recent",
        "teal",
    )
    add_arch_box(
        slide,
        Inches(7.0),
        Inches(4.55),
        Inches(3.4),
        Inches(1.2),
        "Dashboard Angular",
        "Polling 4s, status mix, top MTI, erreurs",
        "green",
    )
    add_footer(slide)

    # Slide 5
    slide = prs.slides.add_slide(layout)
    set_background(slide, COLORS["bg"])
    add_title(slide, "Preuves De Validation", "Resultats verifies pendant cette session")
    add_card(
        slide,
        Inches(0.8),
        Inches(1.9),
        Inches(3.8),
        Inches(1.55),
        "Build backend",
        "Commande: mvn test\nResultat: BUILD SUCCESS\nCouverture constatee: 32 tests verts",
        "green",
    )
    add_card(
        slide,
        Inches(4.8),
        Inches(1.9),
        Inches(3.8),
        Inches(1.55),
        "Build frontend",
        "Commande: npm run build\nResultat: build Angular OK\nSortie: frontend/dist/frontend",
        "teal",
    )
    add_card(
        slide,
        Inches(8.8),
        Inches(1.9),
        Inches(3.7),
        Inches(1.55),
        "Docker compose",
        "Commande: docker compose config\nResultat: fichier valide\nReserve: cle version obsolete a nettoyer",
        "gold",
    )
    add_card(
        slide,
        Inches(0.8),
        Inches(3.8),
        Inches(5.6),
        Inches(1.7),
        "Empreinte du depot",
        "6 commits entre le 1 mars 2026 et le 31 mars 2026.\n30 fichiers Java de production, 7 fichiers de test, 12 fichiers frontend dans src/.",
        "navy",
    )
    add_card(
        slide,
        Inches(6.65),
        Inches(3.8),
        Inches(5.85),
        Inches(1.7),
        "Maturite de la branche",
        "Le worktree est encore tres actif: refonte des controleurs, ajout Angular, nouvelle facade powerCARD, documentation PFE et artefacts statiques.",
        "coral",
    )
    add_footer(slide)

    # Slide 6
    slide = prs.slides.add_slide(layout)
    set_background(slide, COLORS["bg"])
    add_title(slide, "Travaux Recents", "Jalons visibles dans l'historique et le worktree")
    timeline = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(1.2), Inches(2.42), Inches(10.9), Inches(0.08)
    )
    timeline.fill.solid()
    timeline.fill.fore_color.rgb = COLORS["line"]
    timeline.line.fill.background()
    add_timeline_node(
        slide,
        Inches(1.55),
        "Demarrage",
        "01/03",
        ["Initialisation du depot", "Base Spring Boot", "Premiere integration jPOS"],
        "navy",
    )
    add_timeline_node(
        slide,
        Inches(4.55),
        "Stabilisation",
        "09/03",
        ["API REST et Swagger", "Mock switch", "Tests du codec"],
        "green",
    )
    add_timeline_node(
        slide,
        Inches(7.45),
        "Extension",
        "24/03",
        ["Montage Angular", "Dashboard monitoring", "Build frontend + dist"],
        "teal",
    )
    add_timeline_node(
        slide,
        Inches(10.25),
        "Consolidation",
        "31/03",
        ["Facade XML powerCARD", "Refonte controller ISO8583", "Mise a jour docs PFE"],
        "gold",
    )
    add_footer(slide)

    # Slide 7
    slide = prs.slides.add_slide(layout)
    set_background(slide, COLORS["bg"])
    add_title(slide, "Points D'attention", "Ce qui reste a consolider avant une livraison plus propre")
    add_bullets(
        slide,
        Inches(0.95),
        Inches(1.9),
        Inches(11.7),
        Inches(4.9),
        [
            "Documentation en decalage: README, roadmap et architecture ne refletent pas encore toute la surface fonctionnelle actuelle.",
            "Branche non figee: 39 fichiers modifies, 35 nouveaux et 3 supprimes au moment de l'audit.",
            "Monitoring en memoire seulement: pas de persistance, pas d'export Prometheus, pas d'historisation longue.",
            "Canal TCP synchrone et sans mecanismes de resilience avances (retry, circuit breaker, file d'attente).",
            "Securite encore absente dans l'API: pas d'authentification, pas d'autorisation, pas de durcissement d'exposition.",
            "docker-compose.yml comporte une cle version obsolete; la configuration est valide mais merite nettoyage.",
        ],
        18,
        "ink",
    )
    add_footer(slide)

    # Slide 8
    slide = prs.slides.add_slide(layout)
    set_background(slide, COLORS["bg"])
    add_title(slide, "Prochaines Etapes", "Plan court terme recommande")
    add_card(
        slide,
        Inches(0.8),
        Inches(1.9),
        Inches(3.8),
        Inches(2.0),
        "1. Figer la version",
        "Nettoyer le worktree, exclure les artefacts inutiles, aligner README et docs techniques sur le code reel.",
        "navy",
    )
    add_card(
        slide,
        Inches(4.75),
        Inches(1.9),
        Inches(3.8),
        Inches(2.0),
        "2. Valider le flux bout en bout",
        "Rejouer une demo complete gateway + mock switch + dashboard + facade XML, idealement via docker compose.",
        "green",
    )
    add_card(
        slide,
        Inches(8.7),
        Inches(1.9),
        Inches(3.8),
        Inches(2.0),
        "3. Industrialiser",
        "Ajouter securite, persistance du monitoring, pipeline CI/CD, tests end-to-end et supervision exportable.",
        "teal",
    )
    add_text(
        slide,
        Inches(0.9),
        Inches(4.55),
        Inches(11.5),
        Inches(1.0),
        "Message cle pour le jury ou le comite projet: la preuve de faisabilite est convaincante; l'effort restant porte surtout sur la finition et l'industrialisation.",
        20,
        "navy",
        True,
    )
    add_footer(slide)

    prs.save(OUTPUT)


if __name__ == "__main__":
    build_presentation()
    print(OUTPUT)
